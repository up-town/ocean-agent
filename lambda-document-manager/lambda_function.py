import json
import boto3
import os
import traceback
import time
import docx
import base64
import fitz

from io import BytesIO
from urllib import parse
from botocore.config import Config
from PIL import Image
from urllib.parse import unquote_plus
from langchain_aws import BedrockEmbeddings
from langchain_community.vectorstores.opensearch_vector_search import OpenSearchVectorSearch
from langchain_community.docstore.document import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from opensearchpy import OpenSearch
from pptx import Presentation
from multiprocessing import Process, Pipe
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_aws import ChatBedrock
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx.enum.shape import WD_INLINE_SHAPE_TYPE
from pypdf import PdfReader   
from opensearchpy import OpenSearch, RequestsHttpConnection, AWSV4SignerAuth

sqs = boto3.client('sqs')
s3_client = boto3.client('s3')  
s3_bucket = os.environ.get('s3_bucket') # bucket name
s3_prefix = os.environ.get('s3_prefix')

meta_prefix = "metadata/"
enableParallelSummary = os.environ.get('enableParallelSummary')
enalbeParentDocumentRetrival = os.environ.get('enalbeParentDocumentRetrival')

opensearch_account = os.environ.get('opensearch_account')
opensearch_passwd = os.environ.get('opensearch_passwd')
opensearch_url = os.environ.get('opensearch_url')
sqsUrl = os.environ.get('sqsUrl')
doc_prefix = s3_prefix+'/'
LLM_for_chat = json.loads(os.environ.get('LLM_for_chat'))
LLM_for_multimodal= json.loads(os.environ.get('LLM_for_multimodal'))
LLM_embedding = json.loads(os.environ.get('LLM_embedding'))
selected_chat = 0
selected_multimodal = 0
selected_embedding = 0
maxOutputTokens = 4096

roleArn = os.environ.get('roleArn') 
path = os.environ.get('path')
max_object_size = int(os.environ.get('max_object_size'))

supportedFormat = json.loads(os.environ.get('supportedFormat'))
print('supportedFormat: ', supportedFormat)

enableHybridSearch = os.environ.get('enableHybridSearch')
vectorIndexName = os.environ.get('vectorIndexName')

enableImageExtraction = 'true'
enablePageImageExraction = 'true'
pdf_profile = 'ocean'

os_client = OpenSearch(
    hosts = [{
        'host': opensearch_url.replace("https://", ""), 
        'port': 443
    }],
    http_compress = True,
    http_auth=(opensearch_account, opensearch_passwd),
    use_ssl = True,
    verify_certs = True,
    ssl_assert_hostname = False,
    ssl_show_warn = False,
)
    
def delete_document_if_exist(metadata_key):
    try: 
        s3r = boto3.resource("s3")
        bucket = s3r.Bucket(s3_bucket)
        objs = list(bucket.objects.filter(Prefix=metadata_key))
        print('objs: ', objs)
        
        if(len(objs)>0):
            doc = s3r.Object(s3_bucket, metadata_key)
            meta = doc.get()['Body'].read().decode('utf-8')
            print('meta: ', meta)
            
            ids = json.loads(meta)['ids']
            print('ids: ', ids)
            
            # delete ids
            result = vectorstore.delete(ids)
            print('result: ', result)   
            
            # delete files 
            files = json.loads(meta)['files']
            print('files: ', files)
            
            for file in files:
                s3r.Object(s3_bucket, file).delete()
                print('delete file: ', file)
            
        else:
            print('no meta file: ', metadata_key)
            
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to create meta file")

def get_chat():
    global selected_chat
    profile = LLM_for_chat[selected_chat]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'selected_chat: {selected_chat}, bedrock_region: {bedrock_region}, modelId: {modelId}')
                          
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    parameters = {
        "max_tokens":maxOutputTokens,     
        "temperature":0.1,
        "top_k":250,
        "top_p":0.9,
        "stop_sequences": [HUMAN_PROMPT]
    }
    # print('parameters: ', parameters)

    chat = ChatBedrock(   # new chat model
        model_id=modelId,
        client=boto3_bedrock, 
        model_kwargs=parameters,
    )    
    
    selected_chat = selected_chat + 1
    if selected_chat == len(LLM_for_chat):
        selected_chat = 0
    
    return chat

def get_multimodal():
    global selected_multimodal
    
    profile = LLM_for_multimodal[selected_multimodal]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'LLM: {selected_multimodal}, bedrock_region: {bedrock_region}, modelId: {modelId}')
                          
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    parameters = {
        "max_tokens":maxOutputTokens,     
        "temperature":0.1,
        "top_k":250,
        "top_p":0.9,
        "stop_sequences": [HUMAN_PROMPT]
    }
    # print('parameters: ', parameters)

    multimodal = ChatBedrock(   # new chat model
        model_id=modelId,
        client=boto3_bedrock, 
        model_kwargs=parameters,
    )    
    
    selected_multimodal = selected_multimodal + 1
    if selected_multimodal == len(LLM_for_multimodal):
        selected_multimodal = 0
    
    return multimodal

def get_embedding():
    global selected_embedding
    profile = LLM_embedding[selected_embedding]
    bedrock_region =  profile['bedrock_region']
    model_id = profile['model_id']
    print(f'selected_embedding: {selected_embedding}, bedrock_region: {bedrock_region}')
    
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region, 
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    
    bedrock_embedding = BedrockEmbeddings(
        client=boto3_bedrock,
        region_name = bedrock_region,
        model_id = model_id
    )  
    
    selected_embedding = selected_embedding + 1
    if selected_embedding == len(LLM_embedding):
        selected_embedding = 0
    
    return bedrock_embedding

bedrock_embeddings = get_embedding()

index_name = vectorIndexName
vectorstore = OpenSearchVectorSearch(
    index_name=index_name,  
    is_aoss = False,
    #engine="faiss",  # default: nmslib
    embedding_function = bedrock_embeddings,
    opensearch_url = opensearch_url,
    http_auth=(opensearch_account, opensearch_passwd),
)  

def store_document_for_opensearch(file_type, key):
    print('upload to opensearch: ', key) 
    contents, files, tables, subject_company, rating_date = load_document(file_type, key)
    
    if len(contents) == 0:
        print('no contents: ', key)
        return [], files
    
    # contents = str(contents).replace("\n"," ") 
    print('length: ', len(contents))
    
    docs = []
    
    # text        
    docs.append(Document(
        page_content=contents,
        metadata={
            'name': key,
            'url': path+parse.quote(key),
            'subject_company': subject_company,
            'rating_date': rating_date
        }
    ))
        
    # table
    for table in tables:
        docs.append(Document(
            page_content=table['body'],
            metadata={
                'name': table['name'],
                'url': path+parse.quote(table['name']),
                'page': table['page'],
                'subject_company': subject_company,
                'rating_date': rating_date
            }
        ))  
    print('docs: ', docs)

    ids = add_to_opensearch(docs, key)
    
    return ids, files

# ---->
def store_image_for_opensearch(key, page, subject_company, rating_date):
    print('extract text from an image: ', key) 
                                            
    image_obj = s3_client.get_object(Bucket=s3_bucket, Key=key)
                        
    image_content = image_obj['Body'].read()
    img = Image.open(BytesIO(image_content))
                        
    width, height = img.size 
    print(f"(original) width: {width}, height: {height}, size: {width*height}")
    
    pos = key.find('/')
    prefix = key[pos+1:pos+4]
    print('img_prefix: ', prefix)    
    if pdf_profile=='ocean' and prefix == "img_":
        area = (0, 175, width, height-175)
        cropped_img = img.crop(area)
            
        width, height = cropped_img.size 
        print(f"(croped) width: {width}, height: {height}, size: {width*height}")
                
    if width < 100 or height < 100:  # skip small size image
        return []
                
    isResized = False
    while(width*height > 5242880):
        width = int(width/2)
        height = int(height/2)
        isResized = True
        print(f"(resized) width: {width}, height: {height}, size: {width*height}")
           
    try:             
        if isResized:
            img = img.resize((width, height))
                             
        buffer = BytesIO()
        img.save(buffer, format="PNG")
        img_base64 = base64.b64encode(buffer.getvalue()).decode("utf-8")
                                                                
        # extract text from the image
        chat = get_multimodal()
        text = extract_text(chat, img_base64, subject_company)
        extracted_text = text[text.find('<result>')+8:len(text)-9] # remove <result> tag
        #print('extracted_text: ', extracted_text)
        
        summary = summary_image(chat, img_base64, subject_company)
        image_summary = summary[summary.find('<result>')+8:len(summary)-9] # remove <result> tag
        #print('image summary: ', image_summary)
        
        if len(extracted_text) > 30:
            contents = f"[이미지 요약]\n{image_summary}\n\n[추출된 텍스트]\n{extracted_text}"
        else:
            contents = f"[이미지 요약]\n{image_summary}"
        print('image contents: ', contents)

        docs = []        
        if len(contents) > 30:
            docs.append(
                Document(
                    page_content=contents,
                    metadata={
                        'name': key,
                        'url': path+parse.quote(key),
                        'page': page,
                        'subject_company': subject_company,
                        'rating_date': rating_date
                    }
                )
            )         
        print('docs size: ', len(docs))
        
        return add_to_opensearch(docs, key)
    
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                
        #raise Exception ("Not able to summary")  
        
        return []

def is_not_exist(index_name):    
    if os_client.indices.exists(index_name):        
        print('use exist index: ', index_name)    
        return False
    else:
        print('no index: ', index_name)
        return True
    
def create_nori_index():
    index_body={
        'settings':{
            "index.knn": True,
            "index.knn.algo_param.ef_search": 512,
            'analysis': {
                'analyzer': {
                    'my_analyzer': {
                        'char_filter': ['html_strip'], 
                        'tokenizer': 'nori',
                        'filter': ['nori_number','lowercase','trim','my_nori_part_of_speech'],
                        'type': 'custom'
                    }
                },
                'tokenizer': {
                    'nori': {
                        'decompound_mode': 'mixed',
                        'discard_punctuation': 'true',
                        'type': 'nori_tokenizer'
                    }
                },
                "filter": {
                    "my_nori_part_of_speech": {
                        "type": "nori_part_of_speech",
                        "stoptags": [
                                "E", "IC", "J", "MAG", "MAJ",
                                "MM", "SP", "SSC", "SSO", "SC",
                                "SE", "XPN", "XSA", "XSN", "XSV",
                                "UNA", "NA", "VSV"
                        ]
                    }
                }
            },
        },
        'mappings': {
            'properties': {
                'vector_field': {
                    'type': 'knn_vector',
                    'dimension': 1024,
                    'method': {
                        "name": "hnsw",
                        "engine": "faiss",
                        "parameters": {
                            "ef_construction": 512,
                            "m": 16
                        }
                    }                  
                }
            }
        }
    }
    
    if(is_not_exist(index_name)):
        try: # create index
            response = os_client.indices.create(
                index_name,
                body=index_body
            )
            print('index was created with nori plugin:', response)
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)                
            #raise Exception ("Not able to create the index")

if enableHybridSearch == 'true':
    create_nori_index()
    
def add_to_opensearch(docs, key):    
    if len(docs) == 0:
        return []    
    #print('docs[0]: ', docs[0])       
    
    objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
    print('objectName: ', objectName)    
    metadata_key = meta_prefix+objectName+'.metadata.json'
    print('meta file name: ', metadata_key)    
    delete_document_if_exist(metadata_key)
        
    ids = []
    if enalbeParentDocumentRetrival == 'true':
        parent_splitter = RecursiveCharacterTextSplitter(
            chunk_size=2000,
            chunk_overlap=100,
            separators=["\n\n", "\n", ".", " ", ""],
            length_function = len,
        )
        child_splitter = RecursiveCharacterTextSplitter(
            chunk_size=400,
            chunk_overlap=50,
            # separators=["\n\n", "\n", ".", " ", ""],
            length_function = len,
        )

        parent_docs = parent_splitter.split_documents(docs)
        print('len(parent_docs): ', len(parent_docs))
        
        if len(parent_docs):
            # print('parent_docs[0]: ', parent_docs[0])
            # parent_doc_ids = [str(uuid.uuid4()) for _ in parent_docs]
            # print('parent_doc_ids: ', parent_doc_ids)
            for i, doc in enumerate(parent_docs):
                doc.metadata["doc_level"] = "parent"
                print(f"parent_docs[{i}]: {doc}")
                    
            try:        
                parent_doc_ids = vectorstore.add_documents(parent_docs, bulk_size = 10000)
                print('parent_doc_ids: ', parent_doc_ids) 
                print('len(parent_doc_ids): ', len(parent_doc_ids))
                
                child_docs = []
                       
                page = subject_company = rating_date = ""
                for i, doc in enumerate(parent_docs):
                    _id = parent_doc_ids[i]
                    sub_docs = child_splitter.split_documents([doc])
                    
                    page = subject_company = rating_date = ""
                    if "page" in doc.metadata:
                        page = doc.metadata["page"]
                    if "subject_company" in doc.metadata:
                        subject_company = doc.metadata["subject_company"]
                    if "rating_date" in doc.metadata:
                        rating_date = doc.metadata["rating_date"]
                    
                    for _doc in sub_docs:
                        _doc.metadata["parent_doc_id"] = _id
                        _doc.metadata["doc_level"] = "child"
                        _doc.metadata["page"] = page
                        _doc.metadata["subject_company"] = subject_company
                        _doc.metadata["rating_date"] = rating_date
                        
                    child_docs.extend(sub_docs)
                print('child_docs: ', child_docs)
                
                child_doc_ids = vectorstore.add_documents(child_docs, bulk_size = 10000)
                print('child_doc_ids: ', child_doc_ids) 
                print('len(child_doc_ids): ', len(child_doc_ids))
                    
                ids = parent_doc_ids+child_doc_ids
            except Exception:
                err_msg = traceback.format_exc()
                print('error message: ', err_msg)                
                #raise Exception ("Not able to add docs in opensearch")                
    else:
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=100,
            separators=["\n\n", "\n", ".", " ", ""],
            length_function = len,
        ) 
        
        documents = text_splitter.split_documents(docs)
        print('len(documents): ', len(documents))
        if len(documents):
            print('documents[0]: ', documents[0])        
            
        try:        
            ids = vectorstore.add_documents(documents, bulk_size = 10000)
            print('response of adding documents: ', ids)
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)
            #raise Exception ("Not able to add docs in opensearch")    
    
    print('len(ids): ', len(ids))
    return ids

def extract_images_from_pdf(reader, key):
    picture_count = 1
    
    extracted_image_files = []
    print('pages: ', len(reader.pages))
    for i, page in enumerate(reader.pages):
        print('page: ', page)
        if '/ProcSet' in page['/Resources']:
            print('Resources/ProcSet: ', page['/Resources']['/ProcSet'])        
        if '/XObject' in page['/Resources']:
            print(f"Resources/XObject[{i}]: {page['/Resources']['/XObject']}")
        
        for image_file_object in page.images:
            print('image_file_object: ', image_file_object)
            
            img_name = image_file_object.name
            print('img_name: ', img_name)
            
            if img_name in extracted_image_files:
                print('skip....')
                continue
            
            extracted_image_files.append(img_name)
            # print('list: ', extracted_image_files)
            
            ext = img_name.split('.')[-1]            
            contentType = ""
            if ext == 'png':
                contentType = 'image/png'
            elif ext == 'jpg' or ext == 'jpeg':
                contentType = 'image/jpeg'
            elif ext == 'gif':
                contentType = 'image/gif'
            elif ext == 'bmp':
                contentType = 'image/bmp'
            elif ext == 'tiff' or ext == 'tif':
                contentType = 'image/tiff'
            elif ext == 'svg':
                contentType = 'image/svg+xml'
            elif ext == 'webp':
                contentType = 'image/webp'
            elif ext == 'ico':
                contentType = 'image/x-icon'
            elif ext == 'eps':
                contentType = 'image/eps'
            # print('contentType: ', contentType)
            
            if contentType:                
                image_bytes = image_file_object.data

                pixels = BytesIO(image_bytes)
                pixels.seek(0, 0)
                            
                # get path from key
                objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
                folder = s3_prefix+'/files/'+objectName+'/'
                # print('folder: ', folder)
                            
                img_key = folder+img_name
                
                response = s3_client.put_object(
                    Bucket=s3_bucket,
                    Key=img_key,
                    ContentType=contentType,
                    Body=pixels
                )
                print('response: ', response)
                            
                # metadata
                img_meta = {   # not used yet
                    'bucket': s3_bucket,
                    'key': img_key,
                    'url': path+img_key,
                    'ext': 'png',
                    'page': i+1,
                    'original': key
                }
                print('img_meta: ', img_meta)
                            
                picture_count += 1
                    
                extracted_image_files.append(img_key)

    print('extracted_image_files: ', extracted_image_files)    
    return extracted_image_files
        
def extract_images_from_pptx(prs, key):
    picture_count = 1
    
    extracted_image_files = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            print('shape type: ', shape.shape_type)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                # image bytes to PIL Image object
                image_bytes = image.blob
                
                pixels = BytesIO(image_bytes)
                pixels.seek(0, 0)
                        
                # get path from key
                objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
                folder = s3_prefix+'/files/'+objectName+'/'
                print('folder: ', folder)
                        
                fname = 'img_'+key.split('/')[-1].split('.')[0]+f"_{picture_count}"  
                print('fname: ', fname)
                        
                img_key = folder+fname+'.png'
                        
                response = s3_client.put_object(
                    Bucket=s3_bucket,
                    Key=img_key,
                    ContentType='image/png',
                    Body=pixels
                )
                print('response: ', response)
                        
                # metadata
                img_meta = { # not used yet
                    'bucket': s3_bucket,
                    'key': img_key,
                    'url': path+img_key,
                    'ext': 'png',
                    'page': i+1,
                    'original': key
                }
                print('img_meta: ', img_meta)
                        
                picture_count += 1
                
                extracted_image_files.append(img_key)
    
    print('extracted_image_files: ', extracted_image_files)    
    return extracted_image_files

def extract_images_from_docx(doc_contents, key):
    picture_count = 1
    extracted_image_files = []
    
    for inline_shape in doc_contents.inline_shapes:
        #print('inline_shape.type: ', inline_shape.type)                
        if inline_shape.type == WD_INLINE_SHAPE_TYPE.PICTURE:
            rId = inline_shape._inline.graphic.graphicData.pic.blipFill.blip.embed
            print('rId: ', rId)
        
            image_part = doc_contents.part.related_parts[rId]
        
            filename = image_part.filename
            print('filename: ', filename)
        
            bytes_of_image = image_part.image.blob
            pixels = BytesIO(bytes_of_image)
            pixels.seek(0, 0)
                    
            # get path from key
            objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
            folder = s3_prefix+'/files/'+objectName+'/'
            print('folder: ', folder)
            
            fname = 'img_'+key.split('/')[-1].split('.')[0]+f"_{picture_count}"  
            print('fname: ', fname)
                            
            ext = filename.split('.')[-1]            
            contentType = ""
            if ext == 'png':
                contentType = 'image/png'
            elif ext == 'jpg' or ext == 'jpeg':
                contentType = 'image/jpeg'
            elif ext == 'gif':
                contentType = 'image/gif'
            elif ext == 'bmp':
                contentType = 'image/bmp'
            elif ext == 'tiff' or ext == 'tif':
                contentType = 'image/tiff'
            elif ext == 'svg':
                contentType = 'image/svg+xml'
            elif ext == 'webp':
                contentType = 'image/webp'
            elif ext == 'ico':
                contentType = 'image/x-icon'
            elif ext == 'eps':
                contentType = 'image/eps'
            # print('contentType: ', contentType)
                    
            img_key = folder+fname+'.'+ext
            print('img_key: ', img_key)
            
            response = s3_client.put_object(
                Bucket=s3_bucket,
                Key=img_key,
                ContentType=contentType,
                Body=pixels
            )
            print('response: ', response)
                            
            # metadata
            img_meta = { # not used yet
                'bucket': s3_bucket,
                'key': img_key,
                'url': path+img_key,
                'ext': 'png',
                'original': key
            }
            print('img_meta: ', img_meta)
                            
            picture_count += 1
                    
            extracted_image_files.append(img_key)
    
    print('extracted_image_files: ', extracted_image_files)    
    return extracted_image_files

def extract_table_image(page, index, table_count, bbox, key, subject_company, rating_date):
    pixmap_ori = page.get_pixmap()
    # print(f"width: {pixmap_ori.width}, height: {pixmap_ori.height}")
        
    pixmap = page.get_pixmap(dpi=200)  # dpi=300
    #pixels = pixmap.tobytes() # output: jpg
    
    # convert to png
    img = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
    # print(f"width: {pixmap.width}, height: {pixmap.height}")
    
    rate_width = pixmap.width / pixmap_ori.width
    rate_height = pixmap.height / pixmap_ori.height
    # print(f"rate_width={rate_width}, rate_height={rate_height}")
    
    crop_img = img.crop((bbox[0]*rate_width, bbox[1]*rate_height, bbox[2]*rate_width, bbox[3]*rate_height))
    
    pixels = BytesIO()
    crop_img.save(pixels, format='PNG')
    pixels.seek(0, 0)

    # get path from key
    objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
    folder = s3_prefix+'/captures/'+objectName+'/'
                                
    fname = 'table_'+key.split('/')[-1].split('.')[0]+f"_{table_count}"
    
    if subject_company:
        table_meta = {
            "ext": 'png',
            "page": str(index),
            "company": subject_company,
            "date": rating_date
        }
    else:
        table_meta = {
            "ext": 'png',
            "page": str(index)
        }

    response = s3_client.put_object(
        Bucket=s3_bucket,
        Key=folder+fname+'.png',
        ContentType='image/png',
        Metadata = table_meta,
        Body=pixels
    )
    # print('response: ', response)
    
    return folder+fname+'.png'

from pydantic.v1 import BaseModel, Field
def get_profile_of_doc(content: str):
    """Provide profile of document."""
    
    class Profile(BaseModel):
        subject_company: str = Field(description="The value of 'Subject company'")
        rating_date: str = Field(description="The value of 'Rating data'")
    
    subject_company = rating_date = ""
    for attempt in range(5):
        chat = get_chat()
        structured_llm = chat.with_structured_output(Profile, include_raw=True)
    
        info = structured_llm.invoke(content)
        print(f'attempt: {attempt}, info: {info}')
            
        if not info['parsed'] == None:
            parsed_info = info['parsed']
            subject_company = parsed_info.subject_company
            rating_date = parsed_info.rating_date
                            
            print('subject_company: ', subject_company)            
            print('rating_date: ', rating_date)
            break
    return subject_company, rating_date        
                         
# load documents from s3 for pdf and txt
profile_page = ""
def load_document(file_type, key):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, key)
    
    files = []
    tables = []
    contents = ""
    subject_company = rating_date = ""
    if file_type == 'pdf':
        Byte_contents = doc.get()['Body'].read()

        texts = []
        nImages = []
        try: 
            # pdf reader            
            reader = PdfReader(BytesIO(Byte_contents))
            print('pages: ', len(reader.pages))
            
            # extract text
            imgList = []
            for i, page in enumerate(reader.pages):
                print(f"page[{i}]: {page}")
                
                if i<=1 and pdf_profile == 'ocean': # profile page
                    print('skip the first 2 page!')
                    continue
                    
                texts.append(page.extract_text())
                
                # annotation
                #if '/Type' in page:
                #    print(f"Type[{i}]: {page['/Type']}")                
                #if '/Annots' in page:
                #    print(f"Annots[{i}]: {page['/Annots']}")
                #if '/Group' in page:
                #    print(f"Group[{i}]: {page['/Group']}")
                if '/Contents' in page:                
                    print(f"Contents[{i}]: {page['/Contents']}")                    
                #if '/MediaBox' in page:                
                #    print(f"MediaBox[{i}]: {page['/MediaBox']}")                    
                #if '/Parent' in page:
                #    print(f"Parent[{i}]: {page['/Parent']}")
                                
                nImage = 0
                if '/Resources' in page:
                    print(f"Resources[{i}]: {page['/Resources']}")
                    if '/ProcSet' in page['/Resources']:
                        print(f"Resources/ProcSet[{i}]: {page['/Resources']['/ProcSet']}")
                    if '/XObject' in page['/Resources']:
                        print(f"Resources/XObject[{i}]: {page['/Resources']['/XObject']}")                        
                        for j, image in enumerate(page['/Resources']['/XObject']):
                            print(f"image[{j}]: {image}")                                 
                            if image in imgList:
                                print('Duplicated...')
                                continue    
                            else:
                                imgList.append(image)
                                                    
                            Im = page['/Resources']['/XObject'][image]
                            print(f"{image}[{j}]: {Im}")                            
                            nImage = nImage+1
                            
                print(f"# of images of page[{i}] = {nImage}")
                nImages.append(nImage)
                
                # extract metadata
                if pdf_profile == 'ocean' and i==1:
                    print("---> extract metadata from document")
                    print('content: ', texts[i])
                    
                    subject_company, rating_date_ori = get_profile_of_doc(texts[i])
                    print('subject_company: ', subject_company)
                    
                    from datetime import datetime
                    d = datetime.strptime(rating_date_ori, '%d %B %Y')
                    rating_date = str(d)[:10] 
                    print('rating_date: ', rating_date)

            contents = '\n'.join(texts)
                        
            pages = fitz.open(stream=Byte_contents, filetype='pdf')     

            # extract table data
            table_count = 0
            for i, page in enumerate(pages):
                page_tables = page.find_tables()
                
                if page_tables.tables:
                    print('page_tables.tables: ', len(page_tables.tables))

                    for tab in page_tables.tables:    
                        print(tab.to_markdown())    
                        print(f"index: {i}")
                        print(f"bounding box: {tab.bbox}")  # bounding box of the full table
                        #print(f"top-left cell: {tab.cells[0]}")  # top-left cell
                        #print(f"bottom-right cell: {tab.cells[-1]}")  # bottom-right cell
                        print(f"row count: {tab.row_count}, column count: {tab.col_count}") # row and column counts
                        print("\n\n")
                        
                        if tab.row_count>=2:
                            table_image = extract_table_image(page, i, table_count, tab.bbox, key, subject_company, rating_date)
                            table_count += 1
                        
                            tables.append({
                                "body": tab.to_markdown(),
                                "page": str(i),
                                "name": table_image
                            })                    
                            files.append(table_image)

            # extract page images
            if enablePageImageExraction=='true': 
                for i, page in enumerate(pages):
                    print('page: ', page)
                    
                    imgInfo = page.get_image_info()
                    print(f"imgInfo[{i}]: {imgInfo}")         
                    
                    width = height = 0
                    for j, info in enumerate(imgInfo):
                        bbox = info['bbox']
                        print(f"page[{i}] -> bbox[{j}]: {bbox}")
                        if (bbox[2]-bbox[0]>width or bbox[3]-bbox[1]>height) and (bbox[2]-bbox[0]<940 and bbox[3]-bbox[1]<520):
                            width = bbox[2]-bbox[0]
                            height = bbox[3]-bbox[1]
                            print(f"page[{i}] -> (used) width[{j}]: {bbox[2]-bbox[0]}, height[{j}]: {bbox[3]-bbox[1]}")                    
                        print(f"page[{i}] -> (image) width[{j}]: {info['width']}, height[{j}]: {info['height']}")
                        
                    print(f"nImages[{i}]: {nImages[i]}")  # number of XObjects
                    if nImages[i]>=4 or \
                        (nImages[i]>=1 and (width==0 and height==0)) or \
                        (nImages[i]>=1 and (width>=100 or height>=100)):
                        # save current pdf page to image 
                        pixmap = page.get_pixmap(dpi=200)  # dpi=300
                        #pixels = pixmap.tobytes() # output: jpg
                        
                        # convert to png
                        img = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
                        pixels = BytesIO()
                        img.save(pixels, format='PNG')
                        pixels.seek(0, 0)
                                        
                        # get path from key
                        objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
                        folder = s3_prefix+'/captures/'+objectName+'/'
                        print('folder: ', folder)
                                
                        fname = 'img_'+key.split('/')[-1].split('.')[0]+f"_{i}"
                        print('fname: ', fname)          

                        if pdf_profile == 'ocean':
                            img_meta = {
                                "ext": 'png',
                                "page": str(i),
                                "company": subject_company,
                                "date": rating_date
                            }
                        else: 
                            img_meta = {
                                "ext": 'png',
                                "page": str(i)
                            }
                        print('img_meta: ', img_meta)
                               
                        response = s3_client.put_object(
                            Bucket=s3_bucket,
                            Key=folder+fname+'.png',
                            ContentType='image/png',
                            Metadata = img_meta,
                            Body=pixels
                        )
                        print('response: ', response)
                                                        
                        files.append(folder+fname+'.png')
                                    
                contents = '\n'.join(texts)
                
            elif enableImageExtraction == 'true':
                image_files = extract_images_from_pdf(reader, key)
                for img in image_files:
                    files.append(img)
        
        except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to load the pdf file")
                     
    elif file_type == 'pptx':
        Byte_contents = doc.get()['Body'].read()
            
        try:
            prs = Presentation(BytesIO(Byte_contents))

            texts = []
            for i, slide in enumerate(prs.slides):
                text = ""
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = text + shape.text
                texts.append(text)
            contents = '\n'.join(texts)          
            
            if enableImageExtraction == 'true':
                image_files = extract_images_from_pptx(prs, key)                
                for img in image_files:
                    files.append(img)
                    
        except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to load texts from preseation file")
        
    elif file_type == 'docx':
        try:
            Byte_contents = doc.get()['Body'].read()                    
            doc_contents =docx.Document(BytesIO(Byte_contents))

            texts = []
            for i, para in enumerate(doc_contents.paragraphs):
                if(para.text):
                    texts.append(para.text)
                    # print(f"{i}: {para.text}")        
            contents = '\n'.join(texts)            
            # print('contents: ', contents)
            
            # Extract images
            if enableImageExtraction == 'true':
                image_files = extract_images_from_docx(doc_contents, key)                
                for img in image_files:
                    files.append(img)
            
        except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to load docx")   
                
    elif file_type == 'txt' or file_type == 'md':       
        try:  
            contents = doc.get()['Body'].read().decode('utf-8')
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)        
            # raise Exception ("Not able to load the file")
    
    return contents, files, tables, subject_company, rating_date

def isSupported(type):
    for format in supportedFormat:
        if type == format:
            return True
    
    return False
    
def check_supported_type(key, file_type, size):    
    if key.find('/html/') != -1 or key.find('/node_modules/') != -1 or key.find('/.git/') != -1: # do not include html/node_modules folders
        print('html: ', key.find('/html/'))
        return False
    
    if isSupported(file_type):
        if key[0]=='.' or key[key.rfind('/')+1]=='.':
            print(f"Ignore {key} since the filename starts a dot character for macbook.")        
            return False
        elif size > 0 and size<max_object_size:
            return True
    else:
        return False

HUMAN_PROMPT = "\n\nHuman:"
AI_PROMPT = "\n\nAssistant:"
def get_parameter(model_type):
    if model_type=='titan': 
        return {
            "maxTokenCount":1024,
            "stopSequences":[],
            "temperature":0,
            "topP":0.9
        }
    elif model_type=='claude':
        return {
            "max_tokens_to_sample":maxOutputTokens, # 8k    
            "temperature":0.1,
            "top_k":250,
            "top_p":0.9,
            "stop_sequences": [HUMAN_PROMPT]            
        }
        
def extract_text(chat, img_base64, subject_company):
    if subject_company:
        query = f"이 이미지는 {subject_company}에 대한 정보를 포함하고 있습니다. 텍스트를 추출해서 utf8로 변환하세요. <result> tag를 붙여주세요."
    else:
        query = query = "텍스트를 추출해서 utf8로 변환하세요. <result> tag를 붙여주세요."
    
    messages = [
        HumanMessage(
            content=[
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{img_base64}", 
                    },
                },
                {
                    "type": "text", "text": query
                },
            ]
        )
    ]
    
    try: 
        result = chat.invoke(messages)
        
        extracted_text = result.content
        # print('result of text extraction from an image: ', extracted_text)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return extracted_text

def summary_image(chat, img_base64, subject_company):    
    if subject_company:
        query = f"이 이미지는 {subject_company}에 대한 정보를 포함하고 있습니다. 이미지가 의미하는 내용을 풀어서 자세히 알려주세요. <result> tag를 붙여주세요."
    else:
        query = "이미지가 의미하는 내용을 풀어서 자세히 알려주세요. <result> tag를 붙여주세요."
    
    messages = [
        HumanMessage(
            content=[
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{img_base64}", 
                    },
                },
                {
                    "type": "text", "text": query
                },
            ]
        )
    ]
    
    try: 
        result = chat.invoke(messages)
        
        extracted_text = result.content
        # print('summary from an image: ', extracted_text)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return extracted_text

def get_documentId(key, category):
    documentId = category + "-" + key
    documentId = documentId.replace(' ', '_') # remove spaces  
    documentId = documentId.replace(',', '_') # remove commas # not allowed: [ " * \\ < | , > / ? ]
    documentId = documentId.replace('/', '_') # remove slash
    documentId = documentId.lower() # change to lowercase
                
    return documentId

def create_metadata(bucket, key, meta_prefix, s3_prefix, url, category, documentId, ids, files):
    title = key
    timestamp = int(time.time())

    metadata = {
        "Attributes": {
            "_category": category,
            "_source_url": url,
            "_version": str(timestamp),
            "_language_code": "ko"
        },
        "Title": title,
        "DocumentId": documentId,      
        "ids": ids,
        "files": files
    }
    print('metadata: ', metadata)
    
    #objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)]).upper()
    objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
    print('objectName: ', objectName)

    client = boto3.client('s3')
    try: 
        client.put_object(
            Body=json.dumps(metadata), 
            Bucket=bucket, 
            Key=meta_prefix+objectName+'.metadata.json' 
        )
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to create meta file")
    
# load csv documents from s3
def lambda_handler(event, context):
    print('event: ', event)    
    
    documentIds = []
    for record in event['Records']:
        receiptHandle = record['receiptHandle']
        print("receiptHandle: ", receiptHandle)
        
        body = record['body']
        print("body: ", body)
        
        jsonbody = json.loads(body)        
        bucket = jsonbody['bucket']        
        # translate utf8
        key = unquote_plus(jsonbody['key']) # url decoding
        print('bucket: ', bucket)
        print('key: ', key)        
        eventName = jsonbody['type']
        
        start_time = time.time()      
        
        file_type = key[key.rfind('.')+1:len(key)].lower()
        print('file_type: ', file_type)
            
        if eventName == 'ObjectRemoved:Delete':
            if isSupported(file_type):
                objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
                print('objectName: ', objectName)
                
                # get metadata from s3
                metadata_key = meta_prefix+objectName+'.metadata.json'
                print('metadata_key: ', metadata_key)

                documentId = ""
                try: 
                    metadata_obj = s3_client.get_object(Bucket=bucket, Key=metadata_key)
                    metadata_body = metadata_obj['Body'].read().decode('utf-8')
                    metadata = json.loads(metadata_body)
                    print('metadata: ', metadata)
                    documentId = metadata['DocumentId']
                    print('documentId: ', documentId)
                    documentIds.append(documentId)
                except Exception:
                    err_msg = traceback.format_exc()
                    print('err_msg: ', err_msg)
                    # raise Exception ("Not able to get the object")
                    
                if documentId:
                    try: # delete metadata                        
                        delete_document_if_exist(metadata_key)

                        print('delete metadata: ', metadata_key)
                        result = s3_client.delete_object(Bucket=bucket, Key=metadata_key)
                        # print('result of metadata deletion: ', result)

                    except Exception:
                        err_msg = traceback.format_exc()
                        print('err_msg: ', err_msg)
                        # raise Exception ("Not able to delete documents in Kendra")                    
            else: 
                print('This file format is not supported: ', file_type)                
                    
        elif eventName == "ObjectCreated:Put" or eventName == "ObjectCreated:CompleteMultipartUpload":
            size = 0
            page = subject_company = rating_date = ""
            try:
                s3obj = s3_client.get_object(Bucket=bucket, Key=key)
                print(f"Got object: {s3obj}")
                size = int(s3obj['ContentLength'])    
                
                if 'Metadata' in s3obj:
                    if 'page' in s3obj['Metadata']:
                        page = s3obj['Metadata']['page']
                        print('page: ', page)
                    if 'company' in s3obj['Metadata']:
                        subject_company = s3obj['Metadata']['company']
                        print('subject_company: ', subject_company)
                    if 'date' in s3obj['Metadata']:
                        rating_date = s3obj['Metadata']['date']
                        print('rating_date: ', rating_date)
                
                #attributes = ['ETag', 'Checksum', 'ObjectParts', 'StorageClass', 'ObjectSize']
                #result = s3_client.get_object_attributes(Bucket=bucket, Key=key, ObjectAttributes=attributes)  
                #print('result: ', result)            
                #size = int(result['ObjectSize'])
                print('object size: ', size)
            except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to get object info") 
            
            if check_supported_type(key, file_type, size): 
                category = "upload" # for document
                documentId = get_documentId(key, category)                                
                print('documentId: ', documentId)
                
                ids = files = []
                if file_type == 'pdf' or file_type == 'txt' or file_type == 'md' or file_type == 'csv' or file_type == 'pptx' or file_type == 'docx':
                    ids, files = store_document_for_opensearch(file_type, key)   
                                    
                elif file_type == 'png' or file_type == 'jpg' or file_type == 'jpeg':
                    ids = store_image_for_opensearch(key, page, subject_company, rating_date)
                                                                                                         
                create_metadata(bucket=s3_bucket, key=key, meta_prefix=meta_prefix, s3_prefix=s3_prefix, url=path+parse.quote(key), category=category, documentId=documentId, ids=ids, files=files)

            else: # delete if the object is unsupported one for format or size
                try:
                    print('delete the unsupported file: ', key)                                
                    result = s3_client.delete_object(Bucket=bucket, Key=key)
                    print('result of deletion of the unsupported file: ', result)
                            
                except Exception:
                    err_msg = traceback.format_exc()
                    print('err_msg: ', err_msg)
                    # raise Exception ("Not able to delete unsupported file")
                    
        print('processing time: ', str(time.time() - start_time))
        
        # delete queue
        try:
            sqs.delete_message(QueueUrl=sqsUrl, ReceiptHandle=receiptHandle)
        except Exception as e:        
            print('Fail to delete the queue message: ', e)
            
    return {
        'statusCode': 200
    }
